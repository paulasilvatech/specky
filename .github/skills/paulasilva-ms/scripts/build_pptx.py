#!/usr/bin/env python3
"""
build_pptx.py: paulasilva-ms HTML deck -> native, fully editable PowerPoint.

Every element is a native PowerPoint shape: text frames, autoshapes, native
tables, connectors. Nothing is rasterized or flattened. Open the .pptx and
every word, color, number, and box is editable.

The generator does NOT screen-scrape the rendered HTML. It reads two
structured blocks embedded in the deck HTML, which is the single source of
truth for the deck:

  1. <script type="application/json" id="deck-manifest"> ... </script>
     a list of slide specs, each with a `type` (archetype) and `content`
     (a mapping of the archetype's fields to I18N key paths).
  2. const I18N = { ... };  the trilingual content store.

Speaker notes are written into the native PowerPoint notes pane, in the
requested locale. Notes are read from the sibling file
  {deck-stem}_notes_{locale}.md
(produced by export_notes.py when the deck is generated). If that file is
absent, notes fall back to I18N[locale].notes.

Usage:
    python3 build_pptx.py --input deck.html --output deck.pptx --locale en
    python3 build_pptx.py -i deck.html -o deck.pptx --locale pt-BR

--locale is REQUIRED. A PPTX is a single-language deliverable; the slide
content and the speaker notes are both rendered in the locale you name.
"""
import argparse, json, re, sys, math
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError as exc:
    sys.exit(
        f"ERROR: a required dependency is missing ({exc.name}).\n"
        "  pip install -r scripts/requirements.txt\n"
        "  (this script needs python-pptx and lxml)"
    )

# ============================================================
# CANVAS: 16:9 at 13.333" x 7.5" (mirrors the 1600x900 deck stage)
# ============================================================
EMU_IN = 914400
def IN(v): return Emu(int(v * EMU_IN))
PW, PH = 13.333, 7.5
MARGIN_X = 0.70
CHROME_Y = 0.34
CONTENT_TOP = 1.46
CONTENT_BOT = 6.74
CONTENT_W = PW - 2 * MARGIN_X

# ============================================================
# DESIGN TOKENS: mirror the paulasilva-ms deck --ps-color-* set
# ============================================================
INK   = RGBColor(0x1A,0x1A,0x1A)
INK2  = RGBColor(0x3A,0x3A,0x3A)
INK3  = RGBColor(0x73,0x73,0x73)
PAPER = RGBColor(0xFF,0xFF,0xFF)
BG    = RGBColor(0xF7,0xF7,0xF5)
RULE  = RGBColor(0xE5,0xE5,0xE0)
RULE2 = RGBColor(0xCE,0xCE,0xC7)
DARK_BG   = RGBColor(0x14,0x14,0x14)
DARK_SURF = RGBColor(0x1C,0x1C,0x1A)
DARK_INK  = RGBColor(0xF0,0xF0,0xF0)
DARK_INK2 = RGBColor(0xC7,0xC7,0xC2)
DARK_INK3 = RGBColor(0xA8,0xA8,0xA4)
DARK_RULE = RGBColor(0x2E,0x2E,0x2A)

MS = {
    "blue":   {"base": RGBColor(0x00,0xA4,0xEF), "d700": RGBColor(0x00,0x76,0xAC), "l50": RGBColor(0xE5,0xF6,0xFD)},
    "green":  {"base": RGBColor(0x7F,0xBA,0x00), "d700": RGBColor(0x5A,0x85,0x00), "l50": RGBColor(0xF1,0xF8,0xE3)},
    "yellow": {"base": RGBColor(0xFF,0xB9,0x00), "d700": RGBColor(0xB8,0x85,0x00), "l50": RGBColor(0xFF,0xF7,0xE0)},
    "red":    {"base": RGBColor(0xF2,0x50,0x22), "d700": RGBColor(0xB3,0x38,0x16), "l50": RGBColor(0xFF,0xF0,0xEB)},
    "ink":    {"base": INK3, "d700": INK2, "l50": BG},
}
LOGO_ORDER = ["red", "green", "blue", "yellow"]   # 2x2 logo squares

# PPTX-portable fonts: present on every Office install
SANS = "Segoe UI"
MONO = "Consolas"

# identity-locked strings (never translate)
BRAND_TEXT = "Paula Silva | Software Global Black Belt"

# ============================================================
# LOW-LEVEL SHAPE HELPERS  (generic, adapted from the ms-gartner-deck layer)
# ============================================================
def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color

def _noline(shape):
    shape.line.fill.background()

def _line(shape, color, w=1.0):
    shape.line.color.rgb = color; shape.line.width = Pt(w)

def _noshadow(shape):
    spPr = shape._element.spPr
    if spPr.find(qn('a:effectLst')) is None:
        etree.SubElement(spPr, qn('a:effectLst'))

def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False, radius=0.05):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        IN(x), IN(y), IN(w), IN(h))
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    if fill is None: shp.fill.background()
    else: _fill(shp, fill)
    if line is None: _noline(shp)
    else: _line(shp, line, line_w)
    _noshadow(shp)
    return shp

def hline(slide, x, y, w, color, weight=1.0):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x), IN(y), IN(x+w), IN(y))
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    return ln

def text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.0, wrap=True):
    """runs: string, or list of paragraphs; each paragraph a list of (text, opts).
    opts: size, bold, italic, color, font, tracking (pt), caps, align,
    space_before, space_after, line_spacing."""
    box = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        pa = para[0][1] if para and len(para[0]) > 1 else {}
        p.line_spacing = pa.get("line_spacing", line_spacing)
        if pa.get("align"): p.alignment = pa["align"]
        if pa.get("space_before"): p.space_before = Pt(pa["space_before"])
        if pa.get("space_after"): p.space_after = Pt(pa["space_after"])
        for seg in para:
            t, o = (seg[0], seg[1]) if len(seg) > 1 else (seg[0], {})
            r = p.add_run()
            r.text = t.upper() if o.get("caps") else t
            f = r.font
            f.size = Pt(o.get("size", 12))
            f.bold = o.get("bold", False)
            f.italic = o.get("italic", False)
            f.name = o.get("font", SANS)
            f.color.rgb = o.get("color", INK)
            tr = o.get("tracking")
            if tr is not None:
                r._r.get_or_add_rPr().set("spc", str(int(tr * 100)))
    return box

def base_slide(prs, bg=PAPER):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    b = rect(slide, 0, 0, PW, PH, fill=bg)
    spTree = slide.shapes._spTree
    spTree.remove(b._element)
    spTree.insert(2, b._element)
    return slide

def ms_logo(slide, x, y, size=0.20):
    """2x2 Microsoft squares as 4 native rectangles."""
    gap = size * 0.10
    cell = (size - gap) / 2
    for (cx, cy), key in zip([(0,0),(1,0),(0,1),(1,1)], LOGO_ORDER):
        rect(slide, x + cx*(cell+gap), y + cy*(cell+gap), cell, cell, fill=MS[key]["base"])

# ============================================================
# CHROME: brand bar + slide number, light or dark
# ============================================================
def chrome(slide, page_no, total, dark=False):
    ink3 = DARK_INK3 if dark else INK3
    ms_logo(slide, MARGIN_X, CHROME_Y, size=0.20)
    text(slide, MARGIN_X + 0.32, CHROME_Y - 0.02, 7.4, 0.24, [[
        (BRAND_TEXT, {"size": 8.5, "bold": True, "color": ink3, "tracking": 1.6,
                      "caps": True, "font": MONO}),
    ]], anchor=MSO_ANCHOR.MIDDLE)
    text(slide, PW - MARGIN_X - 1.6, CHROME_Y - 0.02, 1.6, 0.24, [[
        (f"{page_no}", {"size": 9, "bold": True, "color": (DARK_INK2 if dark else INK2)}),
        (f"  /  {total}", {"size": 9, "color": ink3}),
    ]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def eyebrow(slide, y, label, accent):
    text(slide, MARGIN_X, y, CONTENT_W, 0.22, [[
        (label, {"size": 10, "bold": True, "color": accent["d700"],
                 "tracking": 2.0, "caps": True, "font": MONO}),
    ]])

def title_block(slide, y, title, *, dark=False, size=30, accent=None):
    """Title that adapts height to estimated wrap. Returns y where content can start."""
    ink = DARK_INK if dark else INK
    chars = len(title)
    cpl = (CONTENT_W * 122.0) / size
    lines = max(1, math.ceil(chars / cpl))
    line_h = size * 1.16 / 72.0
    h = lines * line_h + 0.10
    text(slide, MARGIN_X, y, CONTENT_W, h + 0.2, [[(title, {"size": size, "color": ink})]],
         line_spacing=1.12)
    return y + h + 0.18

def caption(slide, y, txt, *, dark=False):
    ink3 = DARK_INK3 if dark else INK3
    text(slide, MARGIN_X, y, CONTENT_W, 0.6,
         [[(txt, {"size": 11.5, "italic": True, "color": ink3})]], line_spacing=1.35)

# ============================================================
# I18N RESOLUTION
# ============================================================
class Ctx:
    """Resolves manifest content references against I18N[locale]."""
    def __init__(self, i18n_locale):
        self.d = i18n_locale
    def k(self, keypath, default=""):
        if keypath is None: return default
        cur = self.d
        for part in str(keypath).split("."):
            if not isinstance(cur, dict) or part not in cur:
                return default
            cur = cur[part]
        return cur if isinstance(cur, str) else default

# ============================================================
# ARCHETYPE RENDERERS
# Each takes (prs, ctx, content, page_no, total) and adds one native slide.
# `content` is the manifest entry's content dict; text fields are I18N
# key paths resolved through ctx.k(), structural fields are literals.
# ============================================================
def r_cover(prs, ctx, c, page, total):
    slide = base_slide(prs, PAPER)
    chrome(slide, page, total, dark=False)
    accent = MS[c.get("accent", "blue")]
    eyebrow(slide, 1.62, ctx.k(c["eyebrow"]), accent)
    # title with accent segments, height measured from estimated wrap
    tsize = 40
    para, full = [], ""
    for i, seg in enumerate(c["title_segs"]):
        col = MS[seg["accent"]]["base"] if seg.get("accent") else INK
        txt = ctx.k(seg["k"])
        if i > 0 and txt and txt[0] not in ".,;:!?)":
            txt = " " + txt
        full += txt
        para.append((txt, {"size": tsize, "color": col}))
    cpl = (CONTENT_W * 122.0) / tsize
    lines = max(1, math.ceil(len(full) / cpl))
    line_h = tsize * 1.16 / 72.0
    title_h = lines * line_h + 0.10
    title_y = 1.92
    text(slide, MARGIN_X, title_y, CONTENT_W, title_h + 0.2, [para], line_spacing=1.10)
    # subtitle flows directly below the measured title
    sub_y = title_y + title_h + 0.14
    text(slide, MARGIN_X, sub_y, CONTENT_W - 1.2, 0.7,
         [[(ctx.k(c["subtitle"]), {"size": 15, "color": INK2})]], line_spacing=1.35)
    # meta grid anchored near the bottom of the content area
    meta = [("labels.author", "Paula Silva"),
            ("labels.role", "Software Global Black Belt"),
            ("labels.duration", ctx.k(c.get("duration"))),
            ("labels.date", c.get("date", ""))]
    hline(slide, MARGIN_X, 5.84, CONTENT_W, RULE, 1.0)
    colw = CONTENT_W / 4
    for i, (lk, val) in enumerate(meta):
        mx = MARGIN_X + i * colw
        text(slide, mx, 5.98, colw - 0.2, 0.22, [[
            (ctx.k(lk), {"size": 8.5, "bold": True, "color": INK3, "tracking": 1.4,
                         "caps": True, "font": MONO})]])
        text(slide, mx, 6.22, colw - 0.2, 0.3, [[(val, {"size": 13, "color": INK})]])
    return slide

def r_divider(prs, ctx, c, page, total):
    slide = base_slide(prs, DARK_BG)
    chrome(slide, page, total, dark=True)
    accent = MS[c.get("accent", "blue")]
    eyebrow(slide, 1.95, ctx.k(c["eyebrow"]), accent)
    text(slide, MARGIN_X, 2.22, 3.0, 1.4,
         [[(c.get("number", ""), {"size": 92, "bold": True, "color": accent["base"]})]])
    text(slide, MARGIN_X, 3.66, CONTENT_W, 1.2,
         [[(ctx.k(c["title"]), {"size": 40, "color": DARK_INK})]], line_spacing=1.1)
    if c.get("subtitle"):
        text(slide, MARGIN_X, 4.86, CONTENT_W - 1.5, 0.8,
             [[(ctx.k(c["subtitle"]), {"size": 15, "color": DARK_INK2})]], line_spacing=1.35)
    return slide

def r_list(prs, ctx, c, page, total):
    slide = base_slide(prs, PAPER)
    chrome(slide, page, total, dark=False)
    accent = MS[c.get("accent", "blue")]
    eyebrow(slide, CONTENT_TOP, ctx.k(c["eyebrow"]), accent)
    cy = title_block(slide, CONTENT_TOP + 0.28, ctx.k(c["title"]), size=28)
    items = c["items"]
    cap = c.get("caption")
    avail = (CONTENT_BOT - cy) - (0.5 if cap else 0.0)
    row_h = min(0.92, avail / max(1, len(items)))
    for i, it in enumerate(items):
        ry = cy + i * row_h
        text(slide, MARGIN_X, ry, 0.9, row_h, [[
            (it.get("num", f"{i+1:02d}"), {"size": 17, "bold": True, "color": accent["base"],
                                          "font": MONO})]], anchor=MSO_ANCHOR.TOP)
        text(slide, MARGIN_X + 0.95, ry, CONTENT_W - 0.95, row_h, [[
            (ctx.k(it["k"]), {"size": 14, "color": INK2})]], line_spacing=1.3,
            anchor=MSO_ANCHOR.TOP)
        if i < len(items) - 1:
            hline(slide, MARGIN_X + 0.95, ry + row_h - 0.06, CONTENT_W - 0.95, RULE, 0.75)
    if cap:
        caption(slide, CONTENT_BOT - 0.42, ctx.k(cap))
    return slide

def r_question_grid(prs, ctx, c, page, total):
    slide = base_slide(prs, PAPER)
    chrome(slide, page, total, dark=False)
    accent = MS[c.get("accent", "red")]
    eyebrow(slide, CONTENT_TOP, ctx.k(c["eyebrow"]), accent)
    cy = title_block(slide, CONTENT_TOP + 0.28, ctx.k(c["title"]), size=24)
    cards = c["cards"]
    cap = c.get("caption")
    cols = 2 if len(cards) <= 4 else 3
    rows = math.ceil(len(cards) / cols)
    gap = 0.26
    cw = (CONTENT_W - gap * (cols - 1)) / cols
    avail = (CONTENT_BOT - cy) - (0.46 if cap else 0.0)
    ch = (avail - gap * (rows - 1)) / rows
    for idx, card in enumerate(cards):
        col = idx % cols; row = idx // cols
        cx = MARGIN_X + col * (cw + gap)
        cyy = cy + row * (ch + gap)
        acc = MS[card.get("accent", "blue")]
        rect(slide, cx, cyy, cw, ch, fill=PAPER, line=RULE, line_w=1.0)
        rect(slide, cx, cyy, cw, 0.07, fill=acc["base"])
        pad = 0.22
        text(slide, cx + pad, cyy + 0.20, cw - 2*pad, 0.22, [[
            (ctx.k(card["label"]), {"size": 8.5, "bold": True, "color": acc["d700"],
                                    "tracking": 1.5, "caps": True, "font": MONO})]])
        text(slide, cx + pad, cyy + 0.44, cw - 2*pad, 0.5, [[
            (ctx.k(card["q"]), {"size": 16, "color": INK})]], line_spacing=1.12)
        text(slide, cx + pad, cyy + 0.94, cw - 2*pad, ch - 1.1, [[
            (ctx.k(card["hint"]), {"size": 11, "color": INK2})]], line_spacing=1.35)
    if cap:
        caption(slide, CONTENT_BOT - 0.40, ctx.k(cap))
    return slide

def r_pillar_grid(prs, ctx, c, page, total):
    slide = base_slide(prs, PAPER)
    chrome(slide, page, total, dark=False)
    accent = MS[c.get("accent", "green")]
    eyebrow(slide, CONTENT_TOP, ctx.k(c["eyebrow"]), accent)
    cy = title_block(slide, CONTENT_TOP + 0.28, ctx.k(c["title"]), size=24)
    pillars = c["pillars"]
    cap = c.get("caption")
    gap = 0.24
    cw = (CONTENT_W - gap * (len(pillars) - 1)) / len(pillars)
    avail = (CONTENT_BOT - cy) - (0.46 if cap else 0.0)
    ch = avail
    for i, p in enumerate(pillars):
        cx = MARGIN_X + i * (cw + gap)
        acc = MS[p.get("accent", "blue")]
        rect(slide, cx, cy, cw, ch, fill=PAPER, line=RULE, line_w=1.0)
        rect(slide, cx, cy, 0.07, ch, fill=acc["base"])
        pad = 0.22
        text(slide, cx + pad, cy + 0.20, cw - 2*pad, 0.2, [[
            (p.get("num", f"{i+1:02d}"), {"size": 8.5, "bold": True, "color": acc["d700"],
                                         "tracking": 1.5, "caps": True, "font": MONO})]])
        text(slide, cx + pad, cy + 0.42, cw - 2*pad, 0.4, [[
            (ctx.k(p["title"]), {"size": 16, "color": INK})]], line_spacing=1.12)
        text(slide, cx + pad, cy + 0.84, cw - 2*pad, ch - 1.0, [[
            (ctx.k(p["body"]), {"size": 11, "color": INK2})]], line_spacing=1.35)
    if cap:
        caption(slide, CONTENT_BOT - 0.40, ctx.k(cap))
    return slide

def r_layer_rows(prs, ctx, c, page, total):
    slide = base_slide(prs, PAPER)
    chrome(slide, page, total, dark=False)
    accent = MS[c.get("accent", "blue")]
    eyebrow(slide, CONTENT_TOP, ctx.k(c["eyebrow"]), accent)
    cy = title_block(slide, CONTENT_TOP + 0.28, ctx.k(c["title"]), size=24)
    rows = c["rows"]
    cap = c.get("caption")
    gap = 0.12
    avail = (CONTENT_BOT - cy) - (0.46 if cap else 0.0)
    rh = (avail - gap * (len(rows) - 1)) / len(rows)
    for i, row in enumerate(rows):
        ry = cy + i * (rh + gap)
        acc = MS[row.get("accent", "blue")]
        rect(slide, MARGIN_X, ry, CONTENT_W, rh, fill=PAPER, line=RULE, line_w=1.0)
        rect(slide, MARGIN_X, ry, 0.08, rh, fill=acc["base"])
        text(slide, MARGIN_X + 0.26, ry + 0.10, 1.7, 0.22, [[
            (row.get("label", ""), {"size": 8.5, "bold": True, "color": acc["d700"],
                                    "tracking": 1.3, "caps": True, "font": MONO})]])
        text(slide, MARGIN_X + 0.26, ry + 0.30, 2.3, rh - 0.36, [[
            (ctx.k(row["name"]), {"size": 15, "color": acc["d700"]})]], line_spacing=1.1)
        text(slide, MARGIN_X + 2.8, ry, CONTENT_W - 3.0, rh, [[
            (ctx.k(row["desc"]), {"size": 11.5, "color": INK2})]],
            line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)
    if cap:
        caption(slide, CONTENT_BOT - 0.40, ctx.k(cap))
    return slide

def r_roadmap_rows(prs, ctx, c, page, total):
    slide = base_slide(prs, PAPER)
    chrome(slide, page, total, dark=False)
    accent = MS[c.get("accent", "green")]
    eyebrow(slide, CONTENT_TOP, ctx.k(c["eyebrow"]), accent)
    cy = title_block(slide, CONTENT_TOP + 0.28, ctx.k(c["title"]), size=24)
    rows = c["rows"]
    cap = c.get("caption")
    gap = 0.10
    avail = (CONTENT_BOT - cy) - (0.46 if cap else 0.0)
    rh = (avail - gap * (len(rows) - 1)) / len(rows)
    for i, row in enumerate(rows):
        ry = cy + i * (rh + gap)
        acc = MS[row.get("accent", "blue")]
        rect(slide, MARGIN_X, ry, CONTENT_W, rh, fill=PAPER, line=RULE, line_w=1.0)
        rect(slide, MARGIN_X, ry, 0.07, rh, fill=acc["base"])
        text(slide, MARGIN_X + 0.24, ry, 1.5, rh, [[
            (ctx.k(row["phase"]), {"size": 12, "bold": True, "color": acc["d700"],
                                   "font": MONO})]], anchor=MSO_ANCHOR.MIDDLE)
        text(slide, MARGIN_X + 1.8, ry, CONTENT_W - 3.5, rh, [[
            (ctx.k(row["text"]), {"size": 12.5, "color": INK})]],
            line_spacing=1.25, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, PW - MARGIN_X - 1.7, ry, 1.55, rh, [[
            (ctx.k(row["duration"]), {"size": 11, "color": INK3, "font": MONO})]],
            anchor=MSO_ANCHOR.MIDDLE)
    if cap:
        caption(slide, CONTENT_BOT - 0.40, ctx.k(cap))
    return slide

def r_chapter_cover(prs, ctx, c, page, total):
    slide = base_slide(prs, PAPER)
    chrome(slide, page, total, dark=False)
    accent = MS[c.get("accent", "blue")]
    rect(slide, MARGIN_X, CONTENT_TOP + 0.2, 0.10, 3.4, fill=accent["base"])
    text(slide, MARGIN_X + 0.4, CONTENT_TOP, 3.0, 1.7, [[
        (c.get("number", "01"), {"size": 120, "bold": True, "color": accent["l50"]})]])
    text(slide, MARGIN_X + 0.4, CONTENT_TOP + 1.7, CONTENT_W - 0.8, 0.24, [[
        (ctx.k(c["tag"]), {"size": 9.5, "bold": True, "color": accent["d700"],
                           "tracking": 1.6, "caps": True, "font": MONO})]])
    text(slide, MARGIN_X + 0.4, CONTENT_TOP + 1.98, CONTENT_W - 0.8, 0.7, [[
        (ctx.k(c["title"]), {"size": 32, "color": INK})]], line_spacing=1.1)
    text(slide, MARGIN_X + 0.4, CONTENT_TOP + 2.74, CONTENT_W - 2.0, 1.6, [[
        (ctx.k(c["sub"]), {"size": 13.5, "color": INK2})]], line_spacing=1.4)
    return slide

def r_two_col(prs, ctx, c, page, total):
    slide = base_slide(prs, PAPER)
    chrome(slide, page, total, dark=False)
    accent = MS[c.get("accent", "yellow")]
    eyebrow(slide, CONTENT_TOP, ctx.k(c["eyebrow"]), accent)
    cy = title_block(slide, CONTENT_TOP + 0.28, ctx.k(c["title"]), size=24)
    cols = c["cols"]
    cap = c.get("caption")
    gap = 0.3
    cw = (CONTENT_W - gap * (len(cols) - 1)) / len(cols)
    avail = (CONTENT_BOT - cy) - (0.46 if cap else 0.0)
    ch = avail
    for i, col in enumerate(cols):
        cx = MARGIN_X + i * (cw + gap)
        acc = MS[col.get("accent", "blue")]
        rect(slide, cx, cy, cw, ch, fill=PAPER, line=RULE, line_w=1.0)
        rect(slide, cx, cy, cw, 0.07, fill=acc["base"])
        pad = 0.26
        text(slide, cx + pad, cy + 0.22, cw - 2*pad, 0.2, [[
            (col.get("label", ""), {"size": 8.5, "bold": True, "color": acc["d700"],
                                    "tracking": 1.4, "caps": True, "font": MONO})]])
        text(slide, cx + pad, cy + 0.44, cw - 2*pad, 0.4, [[
            (ctx.k(col["title"]), {"size": 17, "color": INK})]], line_spacing=1.12)
        text(slide, cx + pad, cy + 0.90, cw - 2*pad, ch - 1.1, [[
            (ctx.k(col["body"]), {"size": 12, "color": INK2})]], line_spacing=1.4)
    if cap:
        caption(slide, CONTENT_BOT - 0.40, ctx.k(cap))
    return slide

def r_final(prs, ctx, c, page, total):
    slide = base_slide(prs, DARK_BG)
    chrome(slide, page, total, dark=True)
    accent = MS[c.get("accent", "blue")]
    eyebrow(slide, 1.7, ctx.k(c["eyebrow"]), accent)
    text(slide, MARGIN_X, 1.98, CONTENT_W, 1.7, [[
        (ctx.k(c["title"]), {"size": 46, "color": DARK_INK})]], line_spacing=1.06)
    text(slide, MARGIN_X, 3.74, CONTENT_W - 1.4, 1.0, [[
        (ctx.k(c["body"]), {"size": 15, "color": DARK_INK2})]], line_spacing=1.4)
    hline(slide, MARGIN_X, 4.92, CONTENT_W, DARK_RULE, 1.0)
    cards = c["contact"]
    colw = CONTENT_W / 3
    for i, card in enumerate(cards):
        mx = MARGIN_X + i * colw
        text(slide, mx, 5.08, colw - 0.3, 0.22, [[
            (ctx.k(card["label"]), {"size": 8.5, "bold": True, "color": DARK_INK3,
                                    "tracking": 1.6, "caps": True, "font": MONO})]])
        for j, line in enumerate(card["lines"]):
            val = line if line.startswith("$LIT$") else ctx.k(line)
            val = val.replace("$LIT$", "")
            text(slide, mx, 5.32 + j*0.30, colw - 0.3, 0.3, [[
                (val, {"size": 13 if j == 0 else 11.5,
                       "color": DARK_INK if j == 0 else DARK_INK2,
                       "bold": j == 0})]])
    return slide

RENDERERS = {
    "cover": r_cover, "divider": r_divider, "list": r_list,
    "question-grid": r_question_grid, "pillar-grid": r_pillar_grid,
    "layer-rows": r_layer_rows, "roadmap-rows": r_roadmap_rows,
    "chapter-cover": r_chapter_cover, "two-col": r_two_col, "final": r_final,
}

# ============================================================
# DECK HTML PARSING
# ============================================================
def extract_blocks(html):
    """Pull the deck-manifest JSON and the const I18N object out of the HTML."""
    m = re.search(r'<script[^>]*id="deck-manifest"[^>]*>(.*?)</script>', html, re.S)
    if not m:
        raise SystemExit("ERROR: no <script id=\"deck-manifest\"> found in the deck HTML.\n"
                         "       A PPTX-exportable deck must embed a slide manifest. See references/deck.md.")
    manifest = json.loads(m.group(1).strip())
    m2 = re.search(r'const I18N = (\{.*?\n\});', html, re.S)
    if not m2:
        raise SystemExit("ERROR: no `const I18N = {...}` found in the deck HTML.")
    i18n = json.loads(m2.group(1))
    return manifest, i18n

def load_notes_md(deck_path, locale):
    """Return {slide_index(1-based): note_text} from the sibling notes md file."""
    stem = Path(deck_path).stem
    # deck stem may end with a locale tag; notes file is {base}_notes_{locale}.md
    base = re.sub(r'_(en|pt-br|es|multi)$', '', stem, flags=re.I)
    cand = Path(deck_path).parent / f"{base}_notes_{locale.lower()}.md"
    if not cand.exists():
        return None
    notes = {}
    cur = None; buf = []
    for line in cand.read_text(encoding="utf-8").splitlines():
        hm = re.match(r'#{1,3}\s*Slide\s+(\d+)', line.strip(), re.I)
        if hm:
            if cur is not None and buf:
                notes[cur] = "\n".join(buf).strip()
            cur = int(hm.group(1)); buf = []
        elif cur is not None:
            buf.append(line)
    if cur is not None and buf:
        notes[cur] = "\n".join(buf).strip()
    return notes

# ============================================================
# BUILD
# ============================================================
def build(input_html, output_pptx, locale):
    html = Path(input_html).read_text(encoding="utf-8")
    manifest, i18n = extract_blocks(html)
    if locale not in i18n:
        raise SystemExit(f"ERROR: locale '{locale}' not in the deck I18N. "
                         f"Available: {', '.join(i18n.keys())}")
    ctx = Ctx(i18n[locale])

    # speaker notes: prefer the sibling .md, fall back to I18N[locale].notes
    md_notes = load_notes_md(input_html, locale)
    i18n_notes = i18n[locale].get("notes", {})
    notes_source = "md file" if md_notes is not None else "I18N fallback"

    prs = Presentation()
    prs.slide_width = IN(PW)
    prs.slide_height = IN(PH)
    total = len(manifest)

    for idx, spec in enumerate(manifest):
        page = idx + 1
        rtype = spec.get("type")
        renderer = RENDERERS.get(rtype)
        if renderer is None:
            raise SystemExit(f"ERROR: slide {page} has unknown archetype '{rtype}'. "
                             f"Known: {', '.join(sorted(RENDERERS))}")
        slide = renderer(prs, ctx, spec["content"], page, total)
        # native speaker notes, in the requested locale
        note = None
        if md_notes is not None:
            note = md_notes.get(page)
        if not note:
            note = i18n_notes.get(f"s{page}")
        if note:
            slide.notes_slide.notes_text_frame.text = note

    prs.save(output_pptx)
    n_notes = sum(1 for i in range(1, total+1)
                  if (md_notes or {}).get(i) or i18n_notes.get(f"s{i}"))
    _verify_output(output_pptx, total)
    print(f"built: {output_pptx}")
    print(f"  locale: {locale}   slides: {total}   speaker notes: {n_notes} ({notes_source})")

def _verify_output(output_pptx, total):
    """Fail loudly if the PPTX was not written, is empty, or has no slides."""
    out_path = Path(output_pptx)
    if not out_path.is_file() or out_path.stat().st_size == 0:
        sys.exit(f"ERROR: output PPTX was not written or is empty: {output_pptx}")
    if total == 0:
        sys.exit("ERROR: no slides were generated; check the input HTML and locale")
    print(f"  size: {out_path.stat().st_size} bytes")

def main():
    ap = argparse.ArgumentParser(description="paulasilva-ms HTML deck -> native editable PPTX.")
    ap.add_argument("--input", "-i", required=True, help="Path to the deck HTML file")
    ap.add_argument("--output", "-o", required=True, help="Path for the output PPTX file")
    ap.add_argument("--locale", "-l", required=True, choices=["en", "pt-BR", "es"],
                    help="REQUIRED. Locale for slide content AND speaker notes.")
    args = ap.parse_args()
    build(args.input, args.output, args.locale)

if __name__ == "__main__":
    main()
